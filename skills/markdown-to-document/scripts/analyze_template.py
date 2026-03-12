"""
Analyze a document template and extract design patterns as JSON.

Usage:
    python analyze_template.py presentation.pptx
    python analyze_template.py presentation.pptx -o style_profile.json
    python analyze_template.py document.docx
    python analyze_template.py spreadsheet.xlsx
    python analyze_template.py project.mpp
    python analyze_template.py project.xml        # MSPDI format

Outputs a JSON style profile describing fonts, colors, layouts, and patterns.
The agent uses this profile to generate new documents matching the original style.

Requires:
    pip install python-pptx (for .pptx), docxtpl (for .docx), openpyxl (for .xlsx)
    pip install mpxj jpype1 + Java 11+ (for .mpp/.mspdi/.mpx)
"""

from __future__ import annotations

import argparse
import json
import os
import re
import subprocess
import sys
from pathlib import Path


def _ensure_package(import_name: str, pip_name: str | None = None) -> bool:
    """Try to import a package; auto-install if non-interactive. Returns True if available."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
        return True
    except (ImportError, OSError):
        pass

    if sys.stdin.isatty():
        answer = input(f'Package "{pip_name}" is not installed. Install it now? [Y/n] ').strip().lower()
        if answer in ("", "y", "yes"):
            subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
            return True
        return False
    else:
        print(f"Auto-installing missing package: {pip_name}", file=sys.stderr)
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
        __import__(import_name)  # verify
        return True


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _color_to_hex(color) -> str | None:
    """Convert a python-pptx / openpyxl color to hex string, or None."""
    if color is None:
        return None
    # python-pptx RGBColor
    if hasattr(color, "rgb"):
        try:
            return f"#{color.rgb}" if color.rgb else None
        except (AttributeError, TypeError):
            return None
    # openpyxl Color
    if hasattr(color, "value") and color.value and color.value != "00000000":
        val = color.value
        if len(val) == 8:
            val = val[2:]  # strip alpha
        return f"#{val}"
    return None


def _emu_to_inches(emu) -> float | None:
    if emu is None:
        return None
    return round(emu / 914400, 2)


def _pt_value(size) -> float | None:
    if size is None:
        return None
    return round(size / 12700, 1)


# ---------------------------------------------------------------------------
# PPTX Analysis
# ---------------------------------------------------------------------------

def analyze_pptx(path: Path) -> dict:
    if not _ensure_package("pptx", "python-pptx"):
        print("Error: python-pptx is required.", file=sys.stderr)
        sys.exit(1)

    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(path))

    # Slide dimensions
    width_in = _emu_to_inches(prs.slide_width)
    height_in = _emu_to_inches(prs.slide_height)

    # Collect per-layout usage and placeholders
    layout_usage: dict[str, dict] = {}
    all_fonts: set[str] = set()
    all_colors: set[str] = set()
    bg_colors: set[str] = set()

    title_fonts: list[dict] = []
    body_fonts: list[dict] = []
    bullet_info: dict[int, dict] = {}

    total_bullets = 0
    total_text_len = 0
    has_notes = False
    has_images = False
    slide_count = len(prs.slides)

    for slide in prs.slides:
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        if layout_name not in layout_usage:
            # Extract placeholder info from layout
            phs = []
            for ph in slide.slide_layout.placeholders:
                phs.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type).split("(")[0].strip() if ph.placeholder_format.type else "UNKNOWN",
                    "position": {
                        "left_inches": _emu_to_inches(ph.left),
                        "top_inches": _emu_to_inches(ph.top),
                    },
                    "size": {
                        "width_inches": _emu_to_inches(ph.width),
                        "height_inches": _emu_to_inches(ph.height),
                    },
                })
            layout_usage[layout_name] = {"usage_count": 0, "placeholders": phs}
        layout_usage[layout_name]["usage_count"] += 1

        # Background
        bg = slide.background
        if bg and bg.fill and bg.fill.type is not None:
            try:
                fc = bg.fill.fore_color
                if fc and fc.rgb:
                    bg_colors.add(f"#{fc.rgb}")
            except Exception:
                pass

        slide_text_len = 0
        slide_bullet_count = 0

        for shape in slide.shapes:
            # Images
            if shape.shape_type and "PICTURE" in str(shape.shape_type):
                has_images = True

            if not hasattr(shape, "text_frame"):
                continue

            is_title = hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0

            for p in shape.text_frame.paragraphs:
                text = p.text.strip()
                if not text:
                    continue

                slide_text_len += len(text)
                level = p.level or 0
                if level >= 0:
                    slide_bullet_count += 1

                # Collect bullet style info
                if level not in bullet_info:
                    bullet_info[level] = {"indent_pt": None, "font_size_pt": None}

                for run in p.runs:
                    font = run.font
                    fname = font.name
                    fsize = _pt_value(font.size) if font.size else None
                    fbold = font.bold
                    fcolor = _color_to_hex(font.color.rgb if font.color and font.color.type is not None else None)

                    if fname:
                        all_fonts.add(fname)
                    if fcolor:
                        all_colors.add(fcolor)

                    entry = {"name": fname, "size_pt": fsize, "bold": fbold, "color": fcolor}
                    if is_title:
                        title_fonts.append(entry)
                    else:
                        body_fonts.append(entry)

                    # Bullet level font size
                    if fsize and bullet_info[level]["font_size_pt"] is None:
                        bullet_info[level]["font_size_pt"] = fsize

                # Fallback: paragraph-level font (covers theme-inherited fonts)
                pfont = p.font
                pf_name = pfont.name if pfont else None
                pf_size = _pt_value(pfont.size) if pfont and pfont.size else None
                pf_bold = pfont.bold if pfont else None
                pf_color = None
                if pfont and pfont.color and pfont.color.type is not None:
                    pf_color = _color_to_hex(pfont.color.rgb)
                if pf_name:
                    all_fonts.add(pf_name)
                # If no run-level info was found, use paragraph-level
                if not p.runs:
                    entry = {"name": pf_name, "size_pt": pf_size, "bold": pf_bold, "color": pf_color}
                    if is_title:
                        title_fonts.append(entry)
                    else:
                        body_fonts.append(entry)
                    if pf_size and bullet_info[level]["font_size_pt"] is None:
                        bullet_info[level]["font_size_pt"] = pf_size

        total_bullets += slide_bullet_count
        total_text_len += slide_text_len

        # Notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                has_notes = True

    # Compute dominant fonts
    def _dominant(entries: list[dict]) -> dict:
        if not entries:
            return {"name": None, "size_pt": None, "bold": None, "color": None}
        # Pick most common non-None name
        names = [e["name"] for e in entries if e["name"]]
        sizes = [e["size_pt"] for e in entries if e["size_pt"]]
        bolds = [e["bold"] for e in entries if e["bold"] is not None]
        colors = [e["color"] for e in entries if e["color"]]

        from collections import Counter
        name = Counter(names).most_common(1)[0][0] if names else None
        size = Counter(sizes).most_common(1)[0][0] if sizes else None
        bold = Counter(bolds).most_common(1)[0][0] if bolds else None
        color = Counter(colors).most_common(1)[0][0] if colors else None
        return {"name": name, "size_pt": size, "bold": bold, "color": color}

    title_dom = _dominant(title_fonts)
    body_dom = _dominant(body_fonts)

    # Title case pattern detection
    title_texts = []
    for slide in prs.slides:
        if slide.shapes.title:
            title_texts.append(slide.shapes.title.text)
    title_case = "Unknown"
    if title_texts:
        tc_count = sum(1 for t in title_texts if t == t.title())
        uc_count = sum(1 for t in title_texts if t == t.upper())
        if uc_count > len(title_texts) / 2:
            title_case = "UPPER CASE"
        elif tc_count > len(title_texts) / 2:
            title_case = "Title Case"
        else:
            title_case = "Sentence case"

    # Build bullet_styles
    bullet_styles = {}
    for lvl in sorted(bullet_info):
        bullet_styles[f"level_{lvl}"] = {
            "indent_pt": bullet_info[lvl]["indent_pt"],
            "font_size_pt": bullet_info[lvl]["font_size_pt"],
        }

    # Build layouts list
    layouts_list = []
    for name, info in layout_usage.items():
        layouts_list.append({
            "name": name,
            "usage_count": info["usage_count"],
            "placeholders": info["placeholders"],
        })

    # Recommendations
    recommendations = []
    if layouts_list:
        most_used = max(layouts_list, key=lambda x: x["usage_count"])
        recommendations.append(f"Use '{most_used['name']}' layout for content slides (most common in this template)")
    avg_bullets = round(total_bullets / max(slide_count, 1), 1)
    if avg_bullets:
        recommendations.append(f"Keep bullet points to {int(avg_bullets)}-{int(avg_bullets)+1} per slide")
    if title_dom["name"]:
        recommendations.append(f"Use {title_dom['name']} font family for consistency")
    if title_dom["color"]:
        recommendations.append(f"Title color: {title_dom['color']}")

    return {
        "format": "pptx",
        "analyzed_file": path.name,
        "slide_count": slide_count,
        "slide_dimensions": {"width_inches": width_in, "height_inches": height_in},
        "layouts_used": layouts_list,
        "fonts": {
            "title_font": title_dom,
            "body_font": body_dom,
            "all_fonts_used": sorted(all_fonts),
        },
        "colors": {
            "theme_colors": sorted(all_colors),
            "background_colors": sorted(bg_colors) if bg_colors else ["#FFFFFF"],
            "accent_colors": sorted(all_colors - {title_dom.get("color")} - {body_dom.get("color")}),
        },
        "bullet_styles": bullet_styles,
        "content_patterns": {
            "avg_bullets_per_slide": avg_bullets,
            "avg_text_length_per_slide": round(total_text_len / max(slide_count, 1)),
            "has_speaker_notes": has_notes,
            "has_images": has_images,
            "title_case_pattern": title_case,
        },
        "recommendations": recommendations,
    }


# ---------------------------------------------------------------------------
# DOCX Analysis
# ---------------------------------------------------------------------------

def analyze_docx(path: Path) -> dict:
    if not _ensure_package("docx", "python-docx"):
        print("Error: python-docx is required.", file=sys.stderr)
        sys.exit(1)

    from docx import Document
    from docx.shared import Pt

    doc = Document(str(path))

    # Collect styles
    styles_used: set[str] = set()
    all_fonts: set[str] = set()
    style_info: dict[str, dict] = {}

    heading_counts: dict[str, int] = {}
    paragraph_count = 0
    uses_tables = len(doc.tables) > 0
    uses_images = False

    for p in doc.paragraphs:
        sname = p.style.name if p.style else "Normal"
        styles_used.add(sname)

        if sname.startswith("Heading"):
            heading_counts[sname] = heading_counts.get(sname, 0) + 1

        paragraph_count += 1

        # Extract style details
        if sname not in style_info:
            font = p.style.font if p.style else None
            si: dict = {"font": None, "size_pt": None, "bold": None, "color": None}
            if font:
                si["font"] = font.name
                si["size_pt"] = _pt_value(font.size) if font.size else None
                si["bold"] = font.bold
                if font.color and font.color.rgb:
                    si["color"] = f"#{font.color.rgb}"
                if font.name:
                    all_fonts.add(font.name)
            style_info[sname] = si

        # Check runs for font info
        for run in p.runs:
            if run.font.name:
                all_fonts.add(run.font.name)

    # Check for images in runs
    for p in doc.paragraphs:
        for run in p.runs:
            if run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                uses_images = True
                break

    # Section info
    section_count = len(doc.sections)
    has_headers = False
    has_footers = False
    for section in doc.sections:
        if section.header and section.header.paragraphs:
            header_text = "".join(p.text for p in section.header.paragraphs).strip()
            if header_text:
                has_headers = True
        if section.footer and section.footer.paragraphs:
            footer_text = "".join(p.text for p in section.footer.paragraphs).strip()
            if footer_text:
                has_footers = True

    # Estimate sections from headings
    h1_count = heading_counts.get("Heading 1", 0)
    total_sections = max(h1_count, 1)

    # Style output — use common style name mapping
    mapped_styles = {}
    style_map = {"Heading 1": "heading_1", "Heading 2": "heading_2", "Heading 3": "heading_3", "Normal": "normal"}
    for sname, key in style_map.items():
        if sname in style_info:
            mapped_styles[key] = style_info[sname]

    return {
        "format": "docx",
        "analyzed_file": path.name,
        "styles": {
            **mapped_styles,
            "all_styles_used": sorted(styles_used),
        },
        "structure": {
            "section_count": total_sections,
            "has_toc": any("TOC" in s for s in styles_used),
            "has_headers_footers": has_headers or has_footers,
            "paragraph_count": paragraph_count,
        },
        "content_patterns": {
            "avg_paragraphs_per_section": round(paragraph_count / max(total_sections, 1), 1),
            "uses_tables": uses_tables,
            "uses_images": uses_images,
        },
        "fonts": {
            "all_fonts_used": sorted(all_fonts),
        },
    }


# ---------------------------------------------------------------------------
# XLSX Analysis
# ---------------------------------------------------------------------------

def analyze_xlsx(path: Path) -> dict:
    if not _ensure_package("openpyxl"):
        print("Error: openpyxl is required.", file=sys.stderr)
        sys.exit(1)

    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)

    sheets = []
    for ws in wb.worksheets:
        row_count = ws.max_row or 0
        col_count = ws.max_column or 0

        # Header detection — check if first row has bold or fill
        has_header = False
        header_style: dict = {}
        data_style: dict = {}

        if row_count > 0 and col_count > 0:
            first_row_cells = [ws.cell(row=1, column=c) for c in range(1, min(col_count + 1, 20))]
            bold_count = sum(1 for c in first_row_cells if c.font and c.font.bold)
            if bold_count > len(first_row_cells) / 2:
                has_header = True

            # Header style from first cell
            if first_row_cells:
                c = first_row_cells[0]
                fill_color = None
                if c.fill and c.fill.fgColor:
                    fill_color = _color_to_hex(c.fill.fgColor)
                font_color = None
                if c.font and c.font.color:
                    font_color = _color_to_hex(c.font.color)
                header_style = {
                    "font": c.font.name if c.font else None,
                    "size_pt": c.font.size if c.font and c.font.size else None,
                    "bold": c.font.bold if c.font else None,
                    "fill_color": fill_color,
                    "font_color": font_color,
                }

            # Data style from second row first cell
            if row_count > 1:
                c2 = ws.cell(row=2, column=1)
                data_style = {
                    "font": c2.font.name if c2.font else None,
                    "size_pt": c2.font.size if c2.font and c2.font.size else None,
                    "alignment": c2.alignment.horizontal if c2.alignment else None,
                }

        # Column widths
        col_widths = []
        for c in range(1, min(col_count + 1, 20)):
            letter = ws.cell(row=1, column=c).column_letter
            dim = ws.column_dimensions.get(letter)
            if dim and dim.width:
                col_widths.append(round(dim.width, 1))
            else:
                col_widths.append(8.43)  # Default Excel column width

        # Formulas — check a sample of cells
        has_formulas = False
        wb_formulas = load_workbook(str(path), data_only=False)
        ws_f = wb_formulas[ws.title]
        for row in ws_f.iter_rows(max_row=min(row_count, 100), max_col=min(col_count, 20)):
            for cell in row:
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    has_formulas = True
                    break
            if has_formulas:
                break

        # Charts
        has_charts = len(ws._charts) > 0

        # Merged cells
        merged = [str(rng) for rng in ws.merged_cells.ranges]

        sheets.append({
            "name": ws.title,
            "row_count": row_count,
            "col_count": col_count,
            "has_header_row": has_header,
            "header_style": header_style,
            "data_style": data_style,
            "column_widths": col_widths[:col_count],
            "has_formulas": has_formulas,
            "has_charts": has_charts,
            "merged_cells": merged if merged else [],
        })

    return {
        "format": "xlsx",
        "analyzed_file": path.name,
        "sheets": sheets,
    }


# ---------------------------------------------------------------------------
# MPP / MSPDI Analysis (MS Project via MPXJ)
# ---------------------------------------------------------------------------

# JVM bootstrap helpers — shared with generate_project.py


def _check_java(min_version: int = 11) -> str | None:
    """Check if Java is available and meets the minimum version. Returns java path or None."""
    resolver = Path(__file__).parent / "runtime_resolver.py"
    if resolver.is_file():
        try:
            result = subprocess.run(
                [sys.executable, str(resolver), "java",
                 "--min-version", str(min_version)],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                info = json.loads(result.stdout)
                if info.get("found"):
                    return info.get("path")
        except Exception:
            pass

    java_cmd = "java"
    try:
        result = subprocess.run(
            [java_cmd, "-version"],
            capture_output=True, text=True, timeout=10
        )
        if result.stderr or result.stdout:
            return java_cmd
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None


def _resolve_java_home(java_path: str) -> Path | None:
    """Resolve the actual JAVA_HOME from a java executable path."""
    resolved = Path(java_path).resolve()
    if resolved.parent.name.lower() == "bin":
        return resolved.parent.parent
    try:
        result = subprocess.run(
            [java_path, "-XshowSettings:properties", "-version"],
            capture_output=True, text=True, timeout=10,
        )
        for line in (result.stderr or "").splitlines():
            if "java.home" in line and "=" in line:
                return Path(line.split("=", 1)[1].strip())
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None


def _find_jvm_dll(jdk_home: Path) -> str | None:
    """Locate the JVM shared library inside a JDK home directory."""
    if sys.platform == "win32":
        candidates = [
            jdk_home / "bin" / "server" / "jvm.dll",
            jdk_home / "bin" / "client" / "jvm.dll",
        ]
    elif sys.platform == "darwin":
        candidates = [
            jdk_home / "lib" / "server" / "libjvm.dylib",
            jdk_home / "lib" / "client" / "libjvm.dylib",
        ]
    else:
        candidates = [
            jdk_home / "lib" / "server" / "libjvm.so",
            jdk_home / "lib" / "client" / "libjvm.so",
        ]
    for c in candidates:
        if c.is_file():
            return str(c)
    return None


def _bootstrap_jvm() -> None:
    """Ensure JVM is started with MPXJ on classpath."""
    java_path = _check_java(min_version=11)
    if not java_path:
        print("Error: Java 11+ required for .mpp/.mspdi analysis.", file=sys.stderr)
        print("Install a JDK 11+ and ensure 'java' is on your PATH.", file=sys.stderr)
        sys.exit(1)

    jdk_home = _resolve_java_home(java_path)
    jvm_path = _find_jvm_dll(jdk_home) if jdk_home else None
    if jdk_home:
        os.environ["JAVA_HOME"] = str(jdk_home)

    _ensure_package("mpxj")
    _ensure_package("jpype", "jpype1")

    import mpxj  # noqa: F401 — adds MPXJ jars to classpath
    import jpype
    if not jpype.isJVMStarted():
        try:
            if jvm_path:
                jpype.startJVM(jvm_path)
            else:
                jpype.startJVM()
        except Exception as e:
            print(f"Error starting JVM: {e}", file=sys.stderr)
            sys.exit(1)


def _java_date_to_str(dt) -> str | None:
    """Convert a Java LocalDateTime/LocalDate to ISO date string, or None."""
    if dt is None:
        return None
    try:
        return str(dt).split("T")[0]  # "2026-04-01T08:00" → "2026-04-01"
    except Exception:
        return str(dt)


def _java_duration_to_days(dur) -> float | None:
    """Convert an MPXJ Duration to days (float), or None."""
    if dur is None:
        return None
    try:
        val = float(dur.getDuration())
        units = str(dur.getUnits())
        if "HOUR" in units:
            return round(val / 8.0, 2)
        if "WEEK" in units:
            return round(val * 5, 2)
        if "MONTH" in units:
            return round(val * 20, 2)
        if "MINUTE" in units:
            return round(val / 480, 2)
        return round(val, 2)  # assume days
    except Exception:
        return None


def _duration_to_str(dur) -> str | None:
    """Convert an MPXJ Duration to a human string like '5d', or None."""
    if dur is None:
        return None
    try:
        val = dur.getDuration()
        units = str(dur.getUnits())
        if "HOUR" in units:
            return f"{val}h"
        if "WEEK" in units:
            return f"{val}w"
        if "MINUTE" in units:
            return f"{val}m"
        return f"{val}d"
    except Exception:
        return str(dur)


def analyze_mpp(path: Path) -> dict:
    """Analyze an MS Project file (.mpp, .mspdi, .mpx) and return a style profile dict."""
    _bootstrap_jvm()

    from org.mpxj.reader import UniversalProjectReader

    reader = UniversalProjectReader()
    project = reader.read(str(path))
    if project is None:
        print(f"Error: MPXJ could not read '{path}'.", file=sys.stderr)
        sys.exit(1)

    # --- Project properties ---
    props = project.getProjectProperties()
    minutes_per_day = None
    try:
        mpd = props.getMinutesPerDay()
        minutes_per_day = float(mpd.getDuration()) if mpd else None
    except Exception:
        pass
    hours_per_day = round(minutes_per_day / 60, 1) if minutes_per_day else 8.0

    minutes_per_week = None
    try:
        mpw = props.getMinutesPerWeek()
        minutes_per_week = float(mpw.getDuration()) if mpw else None
    except Exception:
        pass

    cal = project.getDefaultCalendar()
    cal_name = str(cal.getName() or "") if cal else ""

    project_properties = {
        "name": str(props.getName() or ""),
        "start_date": _java_date_to_str(props.getStartDate()),
        "finish_date": _java_date_to_str(props.getFinishDate()),
        "calendar_name": cal_name,
        "hours_per_day": hours_per_day,
        "days_per_month": int(props.getDaysPerMonth()) if props.getDaysPerMonth() is not None else 20,
        "minutes_per_week": int(minutes_per_week) if minutes_per_week else 2400,
        "currency_symbol": str(props.getCurrencySymbol() or "$"),
        "author": str(props.getAuthor() or ""),
    }

    # --- Calendar ---
    calendar_info = {"name": "", "working_days": [], "working_hours": {}, "exceptions_count": 0}
    if cal:
        calendar_info["name"] = str(cal.getName() or "")
        from org.mpxj import DayType
        day_names = [
            ("MONDAY", "Monday"), ("TUESDAY", "Tuesday"), ("WEDNESDAY", "Wednesday"),
            ("THURSDAY", "Thursday"), ("FRIDAY", "Friday"), ("SATURDAY", "Saturday"),
            ("SUNDAY", "Sunday"),
        ]
        working_days = []
        for enum_name, display_name in day_names:
            try:
                day_type = DayType.valueOf(enum_name)
                if cal.isWorkingDay(day_type):
                    working_days.append(display_name)
            except Exception:
                pass
        calendar_info["working_days"] = working_days
        calendar_info["working_hours"] = {"start": "08:00", "end": "17:00"}
        try:
            exceptions = cal.getCalendarExceptions()
            calendar_info["exceptions_count"] = len(list(exceptions)) if exceptions else 0
        except Exception:
            pass

    # --- Tasks ---
    all_tasks = list(project.getTasks())
    # Filter out the virtual root task (ID 0) which MPXJ adds
    tasks = [t for t in all_tasks if t.getID() is not None and int(t.getID()) > 0]

    total_tasks = len(tasks)
    summary_tasks = 0
    milestones = 0
    regular_tasks = 0
    max_depth = 0
    durations_days = []
    has_baseline = False
    has_actual = False
    wbs_structure = []
    predecessor_total = 0
    pred_types = {"FS": 0, "SS": 0, "FF": 0, "SF": 0}
    pred_counts = []
    has_lag = False

    # Custom fields detection
    text_fields_used = set()
    flag_fields_used = set()
    number_fields_used = set()
    sample_values: dict[str, list] = {}

    for task in tasks:
        is_summary = bool(task.getSummary())
        is_milestone = bool(task.getMilestone())
        if is_summary:
            summary_tasks += 1
        elif is_milestone:
            milestones += 1
        else:
            regular_tasks += 1

        # WBS depth
        outline = int(task.getOutlineLevel()) if task.getOutlineLevel() is not None else 0
        if outline > max_depth:
            max_depth = outline

        # Duration
        dur_days = _java_duration_to_days(task.getDuration())
        if dur_days is not None and not is_summary:
            durations_days.append(dur_days)

        # Baseline / actuals
        if task.getBaselineStart() is not None:
            has_baseline = True
        if task.getActualStart() is not None:
            has_actual = True

        # WBS structure entry
        resources_on_task = []
        try:
            for ra in (task.getResourceAssignments() or []):
                res = ra.getResource()
                if res and res.getName():
                    resources_on_task.append(str(res.getName()))
        except Exception:
            pass

        wbs_entry: dict = {
            "level": outline,
            "name": str(task.getName() or ""),
        }
        if is_summary:
            try:
                children = task.getChildTasks()
                wbs_entry["children_count"] = len(list(children)) if children else 0
            except Exception:
                wbs_entry["children_count"] = 0
            wbs_entry["is_summary"] = True
        else:
            wbs_entry["is_milestone"] = is_milestone
            wbs_entry["duration"] = _duration_to_str(task.getDuration())
            if resources_on_task:
                wbs_entry["resources"] = resources_on_task
        wbs_structure.append(wbs_entry)

        # Predecessors
        preds = task.getPredecessors()
        if preds:
            pred_list = list(preds)
            pred_counts.append(len(pred_list))
            predecessor_total += len(pred_list)
            for pred in pred_list:
                ptype = str(pred.getType()) if pred.getType() else "FS"
                ptype_short = ptype.replace("FINISH_START", "FS").replace("START_START", "SS") \
                    .replace("FINISH_FINISH", "FF").replace("START_FINISH", "SF")
                if ptype_short in pred_types:
                    pred_types[ptype_short] += 1
                lag = pred.getLag()
                if lag and _java_duration_to_days(lag):
                    lag_val = _java_duration_to_days(lag)
                    if lag_val and abs(lag_val) > 0:
                        has_lag = True
        else:
            pred_counts.append(0)

        # Custom fields (Text1-5, Flag1-5, Number1-5)
        for i in range(1, 6):
            try:
                val = task.getText(i)
                if val:
                    field_name = f"Text{i}"
                    text_fields_used.add(field_name)
                    if field_name not in sample_values:
                        sample_values[field_name] = []
                    sv = str(val)
                    if sv not in sample_values[field_name] and len(sample_values[field_name]) < 5:
                        sample_values[field_name].append(sv)
            except Exception:
                pass
            try:
                val = task.getFlag(i)
                if val is not None and bool(val):
                    flag_fields_used.add(f"Flag{i}")
            except Exception:
                pass
            try:
                val = task.getNumber(i)
                if val is not None:
                    number_fields_used.add(f"Number{i}")
            except Exception:
                pass

    # Duration distribution
    dist = {"1d_or_less": 0, "2d_to_5d": 0, "6d_to_10d": 0, "11d_to_20d": 0, "over_20d": 0}
    for d in durations_days:
        if d <= 1:
            dist["1d_or_less"] += 1
        elif d <= 5:
            dist["2d_to_5d"] += 1
        elif d <= 10:
            dist["6d_to_10d"] += 1
        elif d <= 20:
            dist["11d_to_20d"] += 1
        else:
            dist["over_20d"] += 1

    avg_dur = round(sum(durations_days) / len(durations_days), 1) if durations_days else 0

    task_summary_info = {
        "total_tasks": total_tasks,
        "summary_tasks": summary_tasks,
        "milestones": milestones,
        "regular_tasks": regular_tasks,
        "max_wbs_depth": max_depth,
        "avg_duration_days": avg_dur,
        "duration_distribution": dist,
        "has_baseline": has_baseline,
        "has_actual_dates": has_actual,
    }

    # --- Predecessors summary ---
    avg_preds = round(sum(pred_counts) / len(pred_counts), 1) if pred_counts else 0
    max_preds = max(pred_counts) if pred_counts else 0
    predecessor_patterns = {
        "total_links": predecessor_total,
        "types": pred_types,
        "avg_predecessors_per_task": avg_preds,
        "max_predecessors": max_preds,
        "has_lag": has_lag,
    }

    # --- Resources ---
    all_resources = list(project.getResources() or [])
    # Filter out empty / virtual resource at index 0
    resources = [r for r in all_resources if r.getName()]
    resource_by_type: dict[str, int] = {"work": 0, "material": 0, "cost": 0}
    resource_list = []
    for res in resources:
        rtype = str(res.getType()).lower() if res.getType() else "work"
        if "material" in rtype:
            resource_by_type["material"] += 1
            rtype_label = "material"
        elif "cost" in rtype:
            resource_by_type["cost"] += 1
            rtype_label = "cost"
        else:
            resource_by_type["work"] += 1
            rtype_label = "work"

        max_units = None
        try:
            mu = res.getMaxUnits()
            max_units = int(float(str(mu))) if mu is not None else None
        except Exception:
            pass

        std_rate = ""
        try:
            sr = res.getStandardRate()
            if sr is not None:
                std_rate = str(sr)
        except Exception:
            pass

        resource_list.append({
            "name": str(res.getName()),
            "type": rtype_label,
            "group": str(res.getGroup() or ""),
            "max_units": max_units,
            "standard_rate": std_rate,
        })

    resources_info = {
        "total": len(resources),
        "by_type": resource_by_type,
        "list": resource_list,
    }

    # --- Custom fields ---
    custom_fields = {
        "text_fields_used": sorted(text_fields_used),
        "flag_fields_used": sorted(flag_fields_used),
        "number_fields_used": sorted(number_fields_used),
        "sample_values": sample_values,
    }

    # --- Recommendations ---
    recommendations = []
    if max_depth > 0:
        recommendations.append(
            f"Use {max_depth}-level WBS structure to match the original project hierarchy"
        )
    if avg_dur > 0:
        recommendations.append(
            f"Average task duration is {avg_dur} days — keep new tasks in a similar range"
        )
    if len(resources) > 0:
        recommendations.append(
            f"{len(resources)} resources defined — reuse the resource pool for consistency"
        )
    dominant_type = max(pred_types, key=pred_types.get) if predecessor_total > 0 else None
    if dominant_type and pred_types[dominant_type] > 0:
        pct = round(pred_types[dominant_type] / predecessor_total * 100) if predecessor_total else 0
        recommendations.append(
            f"Predominantly {dominant_type} ({pct}%) dependencies — maintain this pattern"
        )
    if milestones > 0 and total_tasks > 0:
        mpct = round(milestones / total_tasks * 100)
        recommendations.append(
            f"{milestones} milestones ({mpct}% of tasks) — include milestones at phase boundaries"
        )
    if has_baseline:
        recommendations.append("Project has baseline data — set baselines on new projects for tracking")
    if text_fields_used:
        recommendations.append(
            f"Custom text fields in use: {', '.join(sorted(text_fields_used))} — populate for consistency"
        )

    # Determine format from extension
    ext = path.suffix.lower()
    fmt = "mpp" if ext == ".mpp" else ext.lstrip(".")

    return {
        "format": fmt,
        "analyzed_file": path.name,
        "project_properties": project_properties,
        "calendar": calendar_info,
        "task_summary": task_summary_info,
        "wbs_structure": wbs_structure,
        "predecessor_patterns": predecessor_patterns,
        "resources": resources_info,
        "custom_fields": custom_fields,
        "recommendations": recommendations,
    }


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

ANALYZERS = {
    ".pptx": analyze_pptx,
    ".docx": analyze_docx,
    ".xlsx": analyze_xlsx,
    ".mpp": analyze_mpp,
    ".mspdi": analyze_mpp,
    ".mpx": analyze_mpp,
    ".xml": analyze_mpp,
}


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Analyze a document template and extract design patterns as JSON."
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to document file (.pptx, .docx, .xlsx, .mpp, .mspdi, .mpx, .xml)."
    )
    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=None,
        help="Output JSON file path. If not provided, prints to stdout."
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    ext = args.input.suffix.lower()
    analyzer = ANALYZERS.get(ext)
    if not analyzer:
        print(f"Error: Unsupported format '{ext}'. Supported: {', '.join(ANALYZERS.keys())}", file=sys.stderr)
        sys.exit(1)

    profile = analyzer(args.input)
    output_json = json.dumps(profile, indent=2, ensure_ascii=False)

    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(output_json, encoding="utf-8")
        print(f"Style profile saved to: {args.output}", file=sys.stderr)
    else:
        print(output_json)


if __name__ == "__main__":
    main()
