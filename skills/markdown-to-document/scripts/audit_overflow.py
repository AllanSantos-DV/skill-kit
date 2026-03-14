"""
Audit a PowerPoint file for text overflow risks.

Usage:
    python audit_overflow.py output.pptx
    python audit_overflow.py output.pptx --budget 0.9
    python audit_overflow.py output.pptx --json

Analyzes each text box and reports capacity usage.
Exit code 1 if any OVERFLOW found, 0 otherwise.

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path


def _ensure_package(import_name: str, pip_name: str | None = None) -> bool:
    """Try to import a package; auto-install if non-interactive. Returns True if available."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
        return True
    except (ImportError, OSError):
        pass

    if sys.stdin.isatty():
        answer = input(f'Package "{pip_name}" is not installed. Install it now? [Y/n] ').strip().lower()
        if answer in ("", "y", "yes"):
            subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
            return True
        return False
    else:
        print(f"Auto-installing missing package: {pip_name}", file=sys.stderr)
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
        __import__(import_name)
        return True


def _get_font_size_pt(shape) -> float:
    """Extract the dominant font size from a shape in points."""
    from pptx.util import Pt

    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font and paragraph.font.size:
            sizes.append(paragraph.font.size / Pt(1))
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(run.font.size / Pt(1))

    return max(sizes) if sizes else 18.0  # default fallback


def _estimate_capacity(width_inches: float, height_inches: float, font_pt: float) -> float:
    """Estimate character capacity of a text box."""
    line_height_inches = font_pt * 1.2 / 72
    char_width_inches = font_pt * 0.5 / 72

    if line_height_inches <= 0 or char_width_inches <= 0:
        return float("inf")

    lines = height_inches / line_height_inches
    chars_per_line = width_inches / char_width_inches
    return lines * chars_per_line


def _suggest_font_size(width_inches: float, height_inches: float, char_count: int) -> int:
    """Binary search for the largest font size that fits the text."""
    best = 8
    for pt in range(8, 73):
        cap = _estimate_capacity(width_inches, height_inches, pt)
        if cap >= char_count:
            best = pt
        else:
            break
    return best


def audit(input_path: Path, budget: float, as_json: bool) -> int:
    """Audit a PPTX for overflow. Returns exit code."""
    if not _ensure_package("pptx", "python-pptx"):
        print("Error: python-pptx is required.", file=sys.stderr)
        return 1

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(input_path))
    results: list[dict] = []
    has_overflow = False

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text
            char_count = len(text)
            if char_count == 0:
                continue

            width_inches = shape.width / Inches(1)
            height_inches = shape.height / Inches(1)
            font_pt = _get_font_size_pt(shape)
            capacity = _estimate_capacity(width_inches, height_inches, font_pt)

            if capacity <= 0:
                continue

            ratio = char_count / capacity
            shape_name = shape.name or f"Shape {shape.shape_id}"

            if ratio > 1.0:
                status = "OVERFLOW"
                icon = "\U0001f534"
                has_overflow = True
                suggested = _suggest_font_size(width_inches, height_inches, char_count)
                suffix = f" \u2014 try font size {suggested}pt"
            elif ratio > budget:
                status = "AT RISK"
                icon = "\U0001f7e1"
                suffix = ""
            else:
                status = "OK"
                icon = "\u2705"
                suffix = ""

            pct = int(ratio * 100)
            entry = {
                "slide": slide_idx,
                "shape": shape_name,
                "status": status,
                "capacity_pct": pct,
                "font_pt": font_pt,
                "char_count": char_count,
            }
            if status == "OVERFLOW":
                entry["suggested_font_pt"] = suggested

            results.append(entry)

            if not as_json:
                print(f"Slide {slide_idx}, Shape \"{shape_name}\": {icon} {status} ({pct}% capacity){suffix}")

    if as_json:
        print(json.dumps(results, indent=2))

    return 1 if has_overflow else 0


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Audit a PowerPoint file for text overflow risks."
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to the .pptx file to audit."
    )
    parser.add_argument(
        "--budget",
        type=float,
        default=0.8,
        help="Capacity threshold for 'AT RISK' warning (default: 0.8 = 80%%)."
    )
    parser.add_argument(
        "--json",
        action="store_true",
        dest="as_json",
        help="Output results as JSON."
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    sys.exit(audit(args.input, args.budget, args.as_json))


if __name__ == "__main__":
    main()
