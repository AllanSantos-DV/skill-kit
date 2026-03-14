"""
Detect leftover placeholder text in a PowerPoint file.

Usage:
    python check_placeholders.py output.pptx
    python check_placeholders.py output.pptx --strict
    python check_placeholders.py output.pptx --json
    python check_placeholders.py output.pptx --patterns extra_patterns.txt

Exit code 1 if placeholders found, 0 if clean.

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from pathlib import Path

# --- Default patterns -----------------------------------------------------------

_BASE_PATTERNS: list[str] = [
    r"(?i)click to add",
    r"(?i)clique para adicionar",
    r"(?i)lorem ipsum",
    r"(?i)xxxx",
    r"(?i)\[company name\]",
    r"(?i)\[nome da empresa\]",
    r"(?i)insert text here",
    r"(?i)insira texto aqui",
    r"(?i)sample text",
    r"(?i)texto de exemplo",
]

_STRICT_PATTERNS: list[str] = [
    r"(?i)\bcontoso\b",
    r"(?i)\bfabrikam\b",
    r"(?i)\bnorthwind\b",
    r"(?i)\badventure works\b",
    r"(?i)\bexample\.com\b",
    r"(?i)\bcontoso\.com\b",
]


def _ensure_package(import_name: str, pip_name: str | None = None) -> bool:
    """Try to import a package; auto-install if non-interactive. Returns True if available."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
        return True
    except (ImportError, OSError):
        pass

    if sys.stdin.isatty():
        answer = input(f'Package "{pip_name}" is not installed. Install it now? [Y/n] ').strip().lower()
        if answer in ("", "y", "yes"):
            subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
            return True
        return False
    else:
        print(f"Auto-installing missing package: {pip_name}", file=sys.stderr)
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
        __import__(import_name)
        return True


def _compile_patterns(strict: bool, extra_file: Path | None) -> list[re.Pattern]:
    """Build compiled regex list from base + optional strict + optional file."""
    raw = list(_BASE_PATTERNS)
    if strict:
        raw.extend(_STRICT_PATTERNS)
    if extra_file:
        for line in extra_file.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line and not line.startswith("#"):
                raw.append(line)
    return [re.compile(p) for p in raw]


def check(input_path: Path, strict: bool, as_json: bool, patterns_file: Path | None) -> int:
    """Check a PPTX for placeholder text. Returns exit code."""
    if not _ensure_package("pptx", "python-pptx"):
        print("Error: python-pptx is required.", file=sys.stderr)
        return 1

    from pptx import Presentation

    prs = Presentation(str(input_path))
    compiled = _compile_patterns(strict, patterns_file)
    findings: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text
            if not text.strip():
                continue

            shape_name = shape.name or f"Shape {shape.shape_id}"

            for pattern in compiled:
                match = pattern.search(text)
                if match:
                    findings.append({
                        "slide": slide_idx,
                        "shape": shape_name,
                        "matched_text": match.group(),
                        "pattern": pattern.pattern,
                    })
                    if not as_json:
                        print(
                            f"Slide {slide_idx}, Shape \"{shape_name}\": "
                            f"placeholder detected \u2014 \"{match.group()}\""
                        )
                    break  # one match per shape is enough

    if as_json:
        print(json.dumps(findings, indent=2))
    elif not findings:
        print("\u2705 No placeholder text found.")

    return 1 if findings else 0


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Detect leftover placeholder text in a PowerPoint file."
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to the .pptx file to check."
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Also detect demo company names and example URLs (Contoso, Fabrikam, example.com, etc.)."
    )
    parser.add_argument(
        "--json",
        action="store_true",
        dest="as_json",
        help="Output results as JSON."
    )
    parser.add_argument(
        "--patterns",
        type=Path,
        default=None,
        metavar="FILE",
        help="Path to a file with additional regex patterns (one per line)."
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)
    if args.patterns and not args.patterns.is_file():
        print(f"Error: Patterns file not found: {args.patterns}", file=sys.stderr)
        sys.exit(1)

    sys.exit(check(args.input, args.strict, args.as_json, args.patterns))


if __name__ == "__main__":
    main()
