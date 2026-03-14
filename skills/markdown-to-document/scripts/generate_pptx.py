"""
Generate a PowerPoint presentation from JSON or Markdown input.

Usage:
    python generate_pptx.py input.json -o output.pptx
    python generate_pptx.py input.json -t template.pptx -o output.pptx
    python generate_pptx.py input.md -o output.pptx
    python generate_pptx.py input.json -t template.pptx -o output.pptx --mode inject

Modes:
    generate (default): Create new slides from scratch using template layouts.
    inject: Replace text in EXISTING template slides, preserving all XML formatting.

Input formats:
    JSON: {"slides": [{"layout": "...", "title": "...", "body": "...", "notes": "...", "image_path": "..."}]}
    JSON (inject mode): {"slides": [{"slide_index": 0, "title": "...", "body": "..."}]}
    Markdown: ## headings become slide titles, content between headings becomes slide body.

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from pathlib import Path


def _ensure_package(import_name: str, pip_name: str | None = None) -> bool:
    """Try to import a package; auto-install if non-interactive. Returns True if available."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
        return True
    except (ImportError, OSError):
        pass

    if sys.stdin.isatty():
        answer = input(f'Package "{pip_name}" is not installed. Install it now? [Y/n] ').strip().lower()
        if answer in ("", "y", "yes"):
            subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
            return True
        return False
    else:
        print(f"Auto-installing missing package: {pip_name}", file=sys.stderr)
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name])
        __import__(import_name)  # verify
        return True


def parse_markdown_to_slides(md_text: str) -> list[dict]:
    """Parse Markdown into slide data. Each ## heading starts a new slide."""
    slides: list[dict] = []
    current: dict | None = None
    for line in md_text.split("\n"):
        if line.startswith("## "):
            if current:
                current["body"] = current["body"].strip()
                slides.append(current)
            current = {"title": line[3:].strip(), "body": "", "notes": ""}
        elif line.startswith("# ") and not slides and current is None:
            # Top-level heading becomes title slide
            current = {"layout": "Title Slide", "title": line[2:].strip(), "body": "", "notes": ""}
        elif current is not None:
            current["body"] += line + "\n"
    if current:
        current["body"] = current["body"].strip()
        slides.append(current)
    return slides


def load_input(input_path: Path) -> list[dict]:
    """Load slides data from JSON or Markdown file."""
    text = input_path.read_text(encoding="utf-8")
    if input_path.suffix.lower() == ".json":
        data = json.loads(text)
        if "slides" not in data:
            print("Error: JSON must have a 'slides' key.", file=sys.stderr)
            sys.exit(1)
        return data["slides"]
    else:
        # Treat as Markdown
        slides = parse_markdown_to_slides(text)
        if not slides:
            print("Error: No slides found. Use ## headings to separate slides.", file=sys.stderr)
            sys.exit(1)
        return slides


def find_layout(prs, layout_name: str | None, fallback_index: int = 1):
    """Find slide layout by name with fallback to index."""
    if layout_name:
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        # Try case-insensitive partial match
        name_lower = layout_name.lower()
        for layout in prs.slide_layouts:
            if name_lower in layout.name.lower():
                return layout
    if fallback_index < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def _parse_bullet_line(line: str) -> tuple[int, str]:
    """Parse a line and return (indent_level, clean_text)."""
    stripped = line.lstrip()
    indent = len(line) - len(stripped)
    level = min(indent // 2, 2)  # Cap at level 2

    # Remove bullet markers
    for marker in ['• ', '- ', '* ', '→ ']:
        if stripped.startswith(marker):
            return level, stripped[len(marker):]

    # Numbered: "1. ", "2. ", etc. — keep the number as text but apply level
    m = re.match(r'^\d+\.\s+', stripped)
    if m:
        return level, stripped

    # No bullet marker — plain text at detected level
    return level, stripped


def _populate_body(placeholder, body_text: str) -> None:
    """Populate a placeholder with formatted paragraphs from body text."""
    from pptx.util import Pt

    tf = placeholder.text_frame

    # Capture default font size from the placeholder's first paragraph before clearing
    base_size = None
    if tf.paragraphs:
        first_p = tf.paragraphs[0]
        if first_p.font and first_p.font.size:
            base_size = first_p.font.size
        elif first_p.runs:
            for run in first_p.runs:
                if run.font.size:
                    base_size = run.font.size
                    break

    if base_size is None:
        base_size = Pt(18)  # Sensible default for body text

    # Scale factors per level — reduce slightly for deeper nesting
    level_scale = {0: 1.0, 1: 0.889, 2: 0.778}

    tf.clear()

    lines = body_text.split("\n")
    first = True
    for line in lines:
        if not line.strip():
            continue

        level, clean_text = _parse_bullet_line(line)

        if first:
            p = tf.paragraphs[0]  # Reuse existing first paragraph
            first = False
        else:
            p = tf.add_paragraph()

        p.text = clean_text
        p.level = level

        scale = level_scale.get(level, 0.778)
        p.font.size = int(base_size * scale)


def generate(slides: list[dict], template_path: Path | None, output_path: Path) -> None:
    """Generate the PowerPoint file."""
    if not _ensure_package("pptx", "python-pptx"):
        print("Error: python-pptx is required.", file=sys.stderr)
        sys.exit(1)
    from pptx import Presentation
    from pptx.util import Inches

    if template_path:
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    for slide_data in slides:
        layout_name = slide_data.get("layout")
        layout = find_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title and slide_data.get("title"):
            slide.shapes.title.text = slide_data["title"]

        # Body — placeholder index 1 with native bullet formatting
        body_text = slide_data.get("body", "")
        if body_text and len(slide.placeholders) > 1:
            _populate_body(slide.placeholders[1], body_text)

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

        # Image
        image_path = slide_data.get("image_path")
        if image_path and Path(image_path).is_file():
            slide.shapes.add_picture(
                image_path,
                Inches(1), Inches(2),
                width=Inches(4)
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Generated: {output_path} ({len(slides)} slides)")


# ---------------------------------------------------------------------------
# Inject mode — replace text in existing template slides, preserving formatting
# ---------------------------------------------------------------------------

_NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


def _inject_text_preserve_format(shape, new_text: str) -> None:
    """Replace text in a shape while preserving all XML formatting."""
    t_elements = shape._element.findall('.//a:t', _NSMAP)
    if not t_elements:
        return
    t_elements[0].text = new_text
    for t_elem in t_elements[1:]:
        t_elem.text = ''


def _inject_body_preserve_format(shape, body_text: str) -> None:
    """Replace body content preserving paragraph and run formatting."""
    from copy import deepcopy
    from lxml import etree

    paragraphs = shape._element.findall('.//a:p', _NSMAP)
    lines = [line for line in body_text.split('\n') if line.strip()]

    if not paragraphs:
        return

    # Fill existing paragraphs with new lines
    for i, line in enumerate(lines):
        if i < len(paragraphs):
            t_elements = paragraphs[i].findall('.//a:t', _NSMAP)
            _level, clean_text = _parse_bullet_line(line)
            if t_elements:
                t_elements[0].text = clean_text
                for t in t_elements[1:]:
                    t.text = ''
        else:
            # More lines than paragraphs — clone last paragraph and append
            new_p = deepcopy(paragraphs[-1])
            t_elements = new_p.findall('.//a:t', _NSMAP)
            _level, clean_text = _parse_bullet_line(line)
            if t_elements:
                t_elements[0].text = clean_text
                for t in t_elements[1:]:
                    t.text = ''
            # Append after last paragraph's parent (the <a:txBody>)
            paragraphs[-1].getparent().append(new_p)

    # Clear excess paragraphs (fewer lines than paragraphs)
    for i in range(len(lines), len(paragraphs)):
        t_elements = paragraphs[i].findall('.//a:t', _NSMAP)
        for t in t_elements:
            t.text = ''


def inject(slides: list[dict], template_path: Path, output_path: Path) -> None:
    """Inject content into existing template slides, preserving formatting."""
    if not _ensure_package("pptx", "python-pptx"):
        print("Error: python-pptx is required.", file=sys.stderr)
        sys.exit(1)
    from pptx import Presentation

    prs = Presentation(str(template_path))
    template_slides = list(prs.slides)

    # Build index map: slide_index → slide_data
    for seq_idx, slide_data in enumerate(slides):
        target_idx = slide_data.get("slide_index", seq_idx)
        if target_idx < 0 or target_idx >= len(template_slides):
            print(
                f"Warning: slide_index {target_idx} out of range "
                f"(template has {len(template_slides)} slides), skipping.",
                file=sys.stderr,
            )
            continue

        slide = template_slides[target_idx]

        # Title
        title_text = slide_data.get("title")
        if title_text and slide.shapes.title:
            _inject_text_preserve_format(slide.shapes.title, title_text)

        # Body — find the first non-title placeholder or text shape
        body_text = slide_data.get("body")
        if body_text:
            body_shape = None
            for ph in slide.placeholders:
                if ph != slide.shapes.title and ph.has_text_frame:
                    body_shape = ph
                    break
            if body_shape is None:
                # Fallback: first non-title shape with text
                for shape in slide.shapes:
                    if shape != slide.shapes.title and shape.has_text_frame:
                        body_shape = shape
                        break
            if body_shape:
                if '\n' in body_text:
                    _inject_body_preserve_format(body_shape, body_text)
                else:
                    _inject_text_preserve_format(body_shape, body_text)

        # Speaker notes
        notes_text = slide_data.get("notes")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Injected: {output_path} ({len(slides)} slides modified, {len(template_slides)} total)")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint presentation from JSON or Markdown input."
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to JSON or Markdown input file."
    )
    parser.add_argument(
        "-t", "--template",
        type=Path,
        default=None,
        help="Path to a .pptx template file. If not provided, a blank presentation is created."
    )
    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=Path("output.pptx"),
        help="Output file path (default: output.pptx)."
    )
    parser.add_argument(
        "--mode",
        choices=["generate", "inject"],
        default="generate",
        help="Mode: 'generate' creates new slides (default), 'inject' replaces text in existing template slides preserving formatting."
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Error: Input file not found: {args.input}", file=sys.stderr)
        sys.exit(1)
    if args.template and not args.template.is_file():
        print(f"Error: Template file not found: {args.template}", file=sys.stderr)
        sys.exit(1)
    if args.mode == "inject" and not args.template:
        print("Error: --mode inject requires a template (-t/--template).", file=sys.stderr)
        sys.exit(1)

    slides = load_input(args.input)

    if args.mode == "inject":
        inject(slides, args.template, args.output)
    else:
        generate(slides, args.template, args.output)


if __name__ == "__main__":
    main()
